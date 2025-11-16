from openai import OpenAI
import generate_prompts
import json
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
from PIL import Image
import pytesseract
import numpy as np
import os
from snowflake.snowpark import Session

def extract_text_from_pdf(file_path):
    reader = PdfReader(file_path)
    text = ""
    for page in reader.pages:
        text += page.extract_text() + "\n"
    return text

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    slide_content = ""
    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_content = ""
        for i, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                text_frame = shape.text_frame
                slide_content = slide_content + text_frame
    return slide_content

def extract_text_from_docx(file_path):
    doc = Document(file_path)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

def extract_text_from_image(file_path):
    img = Image.open(file_path)
    text = pytesseract.image_to_string(img)
    return text

def chunk_text(text, chunk_size = 500, overlap = 10):
    words = text.split()
    chunks = []
    idnum = 1
    start = 1
    while start < len(words):
        end = min(start + chunk_size, len(words))
        chunk = {"chunk_id" : idnum, "text" : " ".join(words[start:end])}
        start += chunk_size - overlap  # move forward but keep overlap
        idnum += 1
        chunks.append(chunk)
    return chunks

def read_and_chunk(file_path, chunk_size=500):
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == ".pdf":
        text = extract_text_from_pdf(file_path)
    elif ext in [".pptx", ".ppt"]:
        text = extract_text_from_pptx(file_path)
    elif ext in [".docx", ".doc"]:
        text = extract_text_from_docx(file_path)
    elif ext in [".txt"]:
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()
    elif ext in [".png", ".jpg", ".jpeg", ".tiff", ".bmp"]:
        text = extract_text_from_image(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")
   
    chunks = chunk_text(text, chunk_size)
    return chunks

def read(file_path) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        text = extract_text_from_pdf(file_path)
    elif ext in [".pptx", ".ppt"]:
        text = extract_text_from_pptx(file_path)
    elif ext in [".docx", ".doc"]:
        text = extract_text_from_docx(file_path)
    elif ext in [".txt"]:
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()
    elif ext in [".png", ".jpg", ".jpeg", ".tiff", ".bmp"]:
        text = extract_text_from_image(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")
    return text
